from typing import Literal

from dotenv import load_dotenv
from langchain_openai import ChatOpenAI
from langgraph.graph import MessagesState, StateGraph
from langgraph.prebuilt import ToolNode

from composio_langgraph import Action, App, ComposioToolSet
from pathlib import Path
import os

load_dotenv()

GOOGLE_SHEET_ID = '1BxiMVs0XRA5nFMdKvBdBZjgmUUqptlbs74OgvE2upms'

# Initialize the toolset for GitHub
composio_toolset = ComposioToolSet()
tools = composio_toolset.get_tools(
    apps=[
        App.GOOGLESHEETS,
        App.CODEINTERPRETER
    ])
tool_node = ToolNode(tools)

model = ChatOpenAI(temperature=0, streaming=True)
model_with_tools = model.bind_tools(tools)

def call_model(state: MessagesState):
    messages = state["messages"]
    response = model_with_tools.invoke(messages)
    return {"messages": [response]}

def should_continue(state: MessagesState) -> Literal["tools", "__end__"]:
    messages = state["messages"]
    last_message = messages[-1]
    if last_message.tool_calls: # type: ignore
        return "tools"
    return "__end__"


workflow = StateGraph(MessagesState)

# Define the two nodes we will cycle between
workflow.add_node("agent", call_model)
workflow.add_node("tools", tool_node)

workflow.add_edge("__start__", "agent")
workflow.add_conditional_edges(
    "agent",
    should_continue,
)
workflow.add_edge("tools", "agent")

app = workflow.compile()

task = f"""
    Create a ppt on the google sheets: {GOOGLE_SHEET_ID}. 
    Create a sandbox
    First retrieve the sheets content, 
    After that run the code to create graphs from the data.
    then pip install python-pptx using code interpreter.
    Then write the code with python-pptx to create a ppt
    
    Ensure that the ppt is detailed and has proper formatting 
    that makes it look good. The graphs in it should be factual. 
    Try to have a few charts and few table with top data points.
    NOTE: Mostly the user passes small sheets, so try to read the whole sheet at once and not via ranges.

    You can test the code to verify installation :

from pptx import Presentation


def main():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.tetx = "Hello World fromm pptx"
    subtitle.text = "using python-ppts!!!"
    prs.save("test.pptx")


if __name__ == "__main__":
    main()
"""
for chunk in app.stream(
    {
        "messages": [
            (
                "human",
                task
            )
        ]
    },
    stream_mode="values",
):
    chunk["messages"][-1].pretty_print()
